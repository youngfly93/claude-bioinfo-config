#!/usr/bin/env python3
"""
Shirley Wu 风格 PPTX 生成器（中英双语 · 纯极简 · 无涂鸦）。

用法:
    python3 build_deck.py <content.json> <output.pptx>

特性: 双语字体混排 · 拉丁字体嵌入(跨平台) · speaker notes · 双图并排 · 内容 lint · 字体自动安装。
设计真源见 references/style.md，内容规格见 references/content-spec.md。
"""
import sys, os, json, shutil, zipfile, io, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ──────────────────────────── 设计 token（真源） ────────────────────────────
ACCENT  = RGBColor(0xCC, 0x00, 0x00)
INK     = RGBColor(0x22, 0x22, 0x22)
INK2    = RGBColor(0x69, 0x69, 0x69)
INK3    = RGBColor(0x88, 0x88, 0x88)
MUTED   = RGBColor(0xCC, 0xCC, 0xCC)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
WHITE85 = RGBColor(0xD9, 0xD9, 0xD9)

F_TITLE_SANS  = "Quicksand"
F_TITLE_SERIF = "Playfair Display"
F_BODY        = "Atkinson Hyperlegible"
F_CJK         = "PingFang SC"

# 拉丁字体文件（用于嵌入 + 自动安装）。CJK 不嵌入（PingFang 系统字体、体积大）。
FONT_DIR = os.path.join(os.path.dirname(__file__), "..", "assets", "fonts")
FONT_FILES = {
    "Quicksand":             {"regular": "Quicksand-Regular.ttf", "bold": "Quicksand-Bold.ttf"},
    "Playfair Display":      {"regular": "PlayfairDisplay-Regular.ttf", "bold": "PlayfairDisplay-Bold.ttf",
                              "italic": "PlayfairDisplay-Italic.ttf"},
    "Atkinson Hyperlegible": {"regular": "AtkinsonHyperlegible-Regular.ttf", "bold": "AtkinsonHyperlegible-Bold.ttf",
                              "italic": "AtkinsonHyperlegible-Italic.ttf"},
}

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"

# ──────────────────────────── 内容 lint ────────────────────────────
LINT = {"title_max": 24, "zh_title_max": 16, "body_lines": 4, "list_items": 6, "statement_max": 40}
_warnings = []
def lint(slide, i):
    t = slide.get("type"); n = i + 1
    def w(msg): _warnings.append(f"  ⚠ 第{n}页({t}): {msg}")
    for k in ("title_en",):
        if slide.get(k) and len(slide[k]) > LINT["title_max"]:
            w(f"{k} 偏长({len(slide[k])}字符>建议{LINT['title_max']})，标题宜≤6词")
    if slide.get("title_zh") and len(slide["title_zh"]) > LINT["zh_title_max"]:
        w(f"中文标题偏长，宜精简")
    b = slide.get("body")
    if isinstance(b, list) and len(b) > LINT["body_lines"]:
        w(f"正文 {len(b)} 行 > 建议{LINT['body_lines']}，考虑拆页（一页一意）")
    if isinstance(slide.get("items"), list) and len(slide["items"]) > LINT["list_items"]:
        w(f"列表 {len(slide['items'])} 项 > 建议{LINT['list_items']}，考虑拆页")
    for k in ("text_zh", "text_en"):
        if slide.get(k) and len(slide[k]) > LINT["statement_max"]:
            w(f"{k} 偏长，金句宜短")

# ──────────────────────────── 底层助手 ────────────────────────────
def set_run_fonts(run, latin, cjk):
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", latin)
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {}); rPr.append(ea)
    ea.set("typeface", cjk)

def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def _set_spacing(para, pct):
    pPr = para._p.get_or_add_pPr()
    ln = pPr.find(qn("a:lnSpc"))
    if ln is None:
        ln = pPr.makeelement(qn("a:lnSpc"), {}); pPr.insert(0, ln)
    pct_el = ln.find(qn("a:spcPct"))
    if pct_el is None:
        pct_el = ln.makeelement(qn("a:spcPct"), {}); ln.append(pct_el)
    pct_el.set("val", str(int(pct * 1000)))

def _set_letter_spacing(run, pt):
    run._r.get_or_add_rPr().set("spc", str(int(pt * 100)))

def add_text(slide, left, top, width, height, segments, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, cjk=None, line_pct=130, para_space=None, wrap=True):
    cjk = cjk or F_CJK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, runs in enumerate(segments):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _set_spacing(p, line_pct)
        if para_space is not None:
            p.space_after = Pt(para_space)
        for rn in runs:
            r = p.add_run(); r.text = rn["t"]
            set_run_fonts(r, rn.get("f", F_BODY), cjk)
            r.font.size = Pt(rn.get("s", 18))
            r.font.color.rgb = rn.get("c", INK)
            r.font.bold = rn.get("bold", False)
            r.font.italic = rn.get("italic", False)
            if rn.get("spacing"):
                _set_letter_spacing(r, rn["spacing"])
    return tb

def add_side_label(slide, text, dark=False):
    if not text: return
    color = MUTED if not dark else INK3
    w, h = Inches(3.0), Inches(0.4)
    cx, cy = Inches(0.30), EMU_H / 2
    box = slide.shapes.add_textbox(int(cx - w / 2), int(cy - h / 2), w, h)
    box.rotation = 270
    tf = box.text_frame; tf.word_wrap = False
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    set_run_fonts(r, F_BODY, F_CJK)
    r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = color

def add_footer(slide, handle, site, dark=False):
    if not (handle or site): return
    txt = "  |  ".join([x for x in (handle, site) if x])
    color = MUTED if not dark else INK3
    add_text(slide, Inches(0), Inches(7.02), EMU_W, Inches(0.35),
             [[{"t": txt, "f": F_BODY, "s": 10, "c": color}]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def add_notes(slide_obj, notes):
    if notes:
        slide_obj.notes_slide.notes_text_frame.text = notes

def _fit(img_w, img_h, box_w, box_h):
    s = min(box_w / img_w, box_h / img_h)
    return int(img_w * s), int(img_h * s)

def _img_size(path):
    try:
        from PIL import Image
        return Image.open(path).size
    except Exception:
        return 1600, 1000

def _place_image(slide, img, box_left, box_top, box_w, box_h):
    if img and os.path.exists(img):
        iw, ih = _img_size(img)
        w, h = _fit(iw, ih, box_w, box_h)
        left = int(box_left + (box_w - w) / 2); top = int(box_top + (box_h - h) / 2)
        slide.shapes.add_picture(img, left, top, w, h)
        return True
    add_text(slide, box_left, int(box_top + box_h / 2), box_w, Inches(0.5),
             [[{"t": f"[缺图: {img}]", "f": F_BODY, "s": 14, "c": ACCENT}]], align=PP_ALIGN.CENTER)
    return False

# ──────────────────────────── 版式 ────────────────────────────
def slide_cover(prs, s):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark = s.get("dark", False); add_bg(sl, INK if dark else WHITE)
    tcol = WHITE if dark else INK; serif = s.get("serif", True)
    tfont = F_TITLE_SERIF if serif else F_TITLE_SANS
    segs = []
    if s.get("title_en"):
        segs.append([{"t": s["title_en"], "f": tfont, "s": 50, "c": tcol, "bold": not serif, "spacing": 0 if serif else 2}])
    if s.get("title_zh"):
        segs.append([{"t": s["title_zh"], "f": tfont, "s": 30, "c": tcol}])
    add_text(sl, Inches(1.2), Inches(2.4), Inches(10.93), Inches(2.4), segs,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_pct=120, para_space=10)
    sub = []
    if s.get("subtitle"): sub.append([{"t": s["subtitle"], "f": F_BODY, "s": 16, "c": INK2 if not dark else WHITE85}])
    if s.get("author"):   sub.append([{"t": s["author"], "f": F_BODY, "s": 14, "c": ACCENT}])
    if sub:
        add_text(sl, Inches(1.2), Inches(4.9), Inches(10.93), Inches(1.2), sub, align=PP_ALIGN.CENTER, line_pct=130, para_space=4)
    return sl

def slide_section(prs, s):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark = s.get("dark", True); add_bg(sl, INK if dark else WHITE)
    tcol = WHITE if dark else INK; segs = []
    if s.get("number"): segs.append([{"t": s["number"], "f": F_TITLE_SANS, "s": 20, "c": ACCENT, "spacing": 3, "bold": True}])
    if s.get("title_en"): segs.append([{"t": s["title_en"], "f": F_TITLE_SANS, "s": 40, "c": tcol, "spacing": 2}])
    if s.get("title_zh"): segs.append([{"t": s["title_zh"], "f": F_TITLE_SANS, "s": 24, "c": tcol if not dark else WHITE85}])
    add_text(sl, Inches(1.2), Inches(2.6), Inches(10.93), Inches(2.3), segs,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_pct=130, para_space=14)
    add_footer(sl, s.get("_handle"), s.get("_site"), dark)
    return sl

def slide_statement(prs, s):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark = s.get("dark", False); add_bg(sl, INK if dark else WHITE)
    tcol = WHITE if dark else INK; segs = []
    if s.get("text_zh"): segs.append([{"t": s["text_zh"], "f": F_BODY, "s": 30, "c": tcol}])
    if s.get("text_en"): segs.append([{"t": s["text_en"], "f": F_TITLE_SERIF, "s": 22, "c": INK2 if not dark else WHITE85, "italic": True}])
    add_text(sl, Inches(1.6), Inches(2.6), Inches(10.13), Inches(2.3), segs,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_pct=140, para_space=14)
    add_side_label(sl, s.get("side_label"), dark); add_footer(sl, s.get("_handle"), s.get("_site"), dark)
    return sl

def slide_content(prs, s):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark = s.get("dark", False); add_bg(sl, INK if dark else WHITE)
    tcol = WHITE if dark else INK; tline = []
    if s.get("title_en"): tline.append({"t": s["title_en"], "f": F_TITLE_SANS, "s": 26, "c": tcol, "spacing": 3})
    if s.get("title_zh"):
        if tline: tline.append({"t": "  " + s["title_zh"], "f": F_TITLE_SANS, "s": 20, "c": INK3})
        else: tline.append({"t": s["title_zh"], "f": F_TITLE_SANS, "s": 26, "c": tcol, "spacing": 1})
    if tline:
        add_text(sl, Inches(1.2), Inches(1.0), Inches(10.93), Inches(1.0), [tline],
                 align=PP_ALIGN.CENTER if s.get("title_center") else PP_ALIGN.LEFT)
    body = s.get("body", [])
    if isinstance(body, str): body = [body]
    segs = [[{"t": b, "f": F_BODY, "s": s.get("body_size", 20), "c": tcol if not dark else WHITE85}] for b in body]
    if segs:
        add_text(sl, Inches(1.8), Inches(2.6), Inches(9.73), Inches(3.2), segs,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_pct=150, para_space=12)
    add_side_label(sl, s.get("side_label"), dark); add_footer(sl, s.get("_handle"), s.get("_site"), dark)
    return sl

def slide_two_column(prs, s):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark = s.get("dark", False); add_bg(sl, INK if dark else WHITE)
    tcol = WHITE if dark else INK
    if s.get("title"):
        add_text(sl, Inches(1.2), Inches(1.0), Inches(10.93), Inches(0.9),
                 [[{"t": s["title"], "f": F_TITLE_SANS, "s": 26, "c": tcol, "spacing": 2}]], align=PP_ALIGN.CENTER)
    for idx, key in enumerate(("left", "right")):
        col = s.get(key) or {}; muted = col.get("muted", False)
        hc = INK3 if muted else tcol; ic = INK3 if muted else (INK2 if not dark else WHITE85)
        x = Inches(1.4) if idx == 0 else Inches(7.0)
        segs = [[{"t": col.get("head", ""), "f": F_BODY, "s": 18, "c": hc, "bold": True}]]
        for it in col.get("items", []):
            segs.append([{"t": it, "f": F_BODY, "s": 16, "c": ic}])
        add_text(sl, x, Inches(2.5), Inches(4.9), Inches(3.6), segs,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_pct=160, para_space=8)
    add_side_label(sl, s.get("side_label"), dark); add_footer(sl, s.get("_handle"), s.get("_site"), dark)
    return sl

def slide_list(prs, s):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark = s.get("dark", False); add_bg(sl, INK if dark else WHITE)
    tcol = WHITE if dark else INK; tline = []
    if s.get("title_en"): tline.append({"t": s["title_en"], "f": F_TITLE_SANS, "s": 26, "c": tcol, "spacing": 3})
    if s.get("title_zh"): tline.append({"t": ("  " if tline else "") + s["title_zh"], "f": F_TITLE_SANS, "s": 20, "c": INK3})
    if tline:
        add_text(sl, Inches(1.2), Inches(1.0), Inches(10.93), Inches(0.9), [tline], align=PP_ALIGN.CENTER)
    items = s.get("items", []); numbered = s.get("numbered", False); segs = []
    for i, it in enumerate(items):
        prefix = f"{i+1}. " if numbered else "·  "
        segs.append([{"t": prefix, "f": F_BODY, "s": 20, "c": ACCENT},
                     {"t": it, "f": F_BODY, "s": 20, "c": tcol if not dark else WHITE85}])
    if segs:
        add_text(sl, Inches(3.0), Inches(2.5), Inches(7.3), Inches(3.6), segs,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_pct=170, para_space=6)
    add_side_label(sl, s.get("side_label"), dark); add_footer(sl, s.get("_handle"), s.get("_site"), dark)
    return sl

def slide_photo(prs, s):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark = s.get("dark", True); add_bg(sl, INK if dark else WHITE)
    wide = s.get("fit") == "wide"
    bw, bh = (Inches(11.6), Inches(5.0)) if wide else (Inches(9.6), Inches(5.2))
    _place_image(sl, s.get("image"), int((EMU_W - bw) / 2), Inches(1.1), bw, bh)
    if s.get("caption"):
        add_text(sl, Inches(1.2), Inches(6.4), Inches(10.93), Inches(0.5),
                 [[{"t": s["caption"], "f": F_BODY, "s": 13, "c": INK3 if dark else INK2}]], align=PP_ALIGN.CENTER)
    add_side_label(sl, s.get("side_label"), dark); add_footer(sl, s.get("_handle"), s.get("_site"), dark)
    return sl

def slide_photo_pair(prs, s):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark = s.get("dark", False); add_bg(sl, INK if dark else WHITE)
    tcol = WHITE if dark else INK
    if s.get("title"):
        add_text(sl, Inches(1.2), Inches(0.7), Inches(10.93), Inches(0.8),
                 [[{"t": s["title"], "f": F_TITLE_SANS, "s": 24, "c": tcol, "spacing": 2}]], align=PP_ALIGN.CENTER)
    imgs = s.get("images", []); caps = s.get("captions", [])
    boxes = [(Inches(0.7), Inches(5.7)), (Inches(6.93), Inches(5.7))]
    for i in range(min(2, len(imgs))):
        bl, bw = boxes[i]
        _place_image(sl, imgs[i], bl, Inches(1.7), bw, Inches(4.0))
        cap = caps[i] if i < len(caps) else None
        if cap:
            add_text(sl, bl, Inches(5.9), bw, Inches(0.6),
                     [[{"t": cap, "f": F_BODY, "s": 12, "c": INK3 if dark else INK2}]], align=PP_ALIGN.CENTER, line_pct=120)
    add_side_label(sl, s.get("side_label"), dark); add_footer(sl, s.get("_handle"), s.get("_site"), dark)
    return sl

def slide_closing(prs, s):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark = s.get("dark", False); add_bg(sl, INK if dark else WHITE)
    tcol = WHITE if dark else INK; segs = []
    if s.get("title_en"): segs.append([{"t": s["title_en"], "f": F_TITLE_SERIF, "s": 40, "c": tcol, "italic": True}])
    if s.get("title_zh"): segs.append([{"t": s["title_zh"], "f": F_TITLE_SANS, "s": 26, "c": tcol}])
    add_text(sl, Inches(1.2), Inches(2.6), Inches(10.93), Inches(1.8), segs,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_pct=130, para_space=10)
    links = s.get("links", [])
    if links:
        add_text(sl, Inches(1.2), Inches(4.6), Inches(10.93), Inches(1.6),
                 [[{"t": l, "f": F_BODY, "s": 14, "c": ACCENT}] for l in links], align=PP_ALIGN.CENTER, line_pct=150, para_space=4)
    add_footer(sl, s.get("_handle"), s.get("_site"), dark)
    return sl

BUILDERS = {
    "cover": slide_cover, "section": slide_section, "statement": slide_statement, "content": slide_content,
    "two_column": slide_two_column, "list": slide_list, "photo": slide_photo, "photo_pair": slide_photo_pair,
    "closing": slide_closing,
}

# ──────────────────────────── 字体自动安装 ────────────────────────────
def ensure_fonts():
    if not os.path.isdir(FONT_DIR): return
    dst = os.path.expanduser("~/Library/Fonts"); os.makedirs(dst, exist_ok=True)
    installed = []
    for f in os.listdir(FONT_DIR):
        if f.lower().endswith((".ttf", ".otf")) and not os.path.exists(os.path.join(dst, f)):
            try: shutil.copy2(os.path.join(FONT_DIR, f), os.path.join(dst, f)); installed.append(f)
            except Exception: pass
    if installed: print(f"  已安装字体到 ~/Library/Fonts: {', '.join(installed)}")

# ──────────────────────────── 字体嵌入 pptx (OOXML) ────────────────────────────
def embed_fonts(pptx_path):
    """把拉丁静态字体嵌入 pptx，使其在未装字体的电脑(含 Windows)也正确渲染。中文 PingFang 不嵌入。"""
    fonts = []  # (typeface, style, filepath)
    for tf, styles in FONT_FILES.items():
        for style, fn in styles.items():
            p = os.path.join(FONT_DIR, fn)
            if os.path.exists(p): fonts.append((tf, style, p))
    if not fonts: return
    zin = zipfile.ZipFile(pptx_path, "r"); items = {n: zin.read(n) for n in zin.namelist()}; zin.close()

    # 1) 字体二进制 parts
    style_tag = {"regular": "regular", "bold": "bold", "italic": "italic", "boldItalic": "boldItalic"}
    by_typeface = {}
    for idx, (tf, style, p) in enumerate(fonts, 1):
        part = f"ppt/fonts/font{idx}.fntdata"
        with open(p, "rb") as fh: items[part] = fh.read()
        by_typeface.setdefault(tf, {})[style] = (part, f"rIdFont{idx}")

    # 2) [Content_Types].xml 加 fntdata default
    ct = etree.fromstring(items["[Content_Types].xml"])
    if not any(d.get("Extension") == "fntdata" for d in ct.findall(f"{{{NS_CT}}}Default")):
        d = etree.SubElement(ct, f"{{{NS_CT}}}Default")
        d.set("Extension", "fntdata"); d.set("ContentType", "application/x-fontdata")
    items["[Content_Types].xml"] = etree.tostring(ct, xml_declaration=True, encoding="UTF-8", standalone=True)

    # 3) presentation.xml.rels 加 font 关系
    rels_name = "ppt/_rels/presentation.xml.rels"
    rels = etree.fromstring(items[rels_name])
    for tf, styles in by_typeface.items():
        for style, (part, rid) in styles.items():
            r = etree.SubElement(rels, f"{{{NS_REL}}}Relationship")
            r.set("Id", rid); r.set("Type", f"{NS_R}/font"); r.set("Target", part[len('ppt/'):])
    items[rels_name] = etree.tostring(rels, xml_declaration=True, encoding="UTF-8", standalone=True)

    # 4) presentation.xml 加 embedTrueTypeFonts + embeddedFontLst
    pres = etree.fromstring(items["ppt/presentation.xml"])
    pres.set("embedTrueTypeFonts", "1")
    lst = etree.SubElement(pres, f"{{{NS_P}}}embeddedFontLst")
    for tf, styles in by_typeface.items():
        ef = etree.SubElement(lst, f"{{{NS_P}}}embeddedFont")
        fo = etree.SubElement(ef, f"{{{NS_P}}}font"); fo.set("typeface", tf)
        for style, (part, rid) in styles.items():
            se = etree.SubElement(ef, f"{{{NS_P}}}{style_tag[style]}")
            se.set(f"{{{NS_R}}}id", rid)
    # 放到正确 schema 位置：notesSz 之后
    anchor = pres.find(f"{{{NS_P}}}notesSz")
    if anchor is not None:
        pres.remove(lst); anchor.addnext(lst)
    items["ppt/presentation.xml"] = etree.tostring(pres, xml_declaration=True, encoding="UTF-8", standalone=True)

    # 5) 重写 zip
    with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for n, data in items.items(): zout.writestr(n, data)
    print(f"  已嵌入 {len(fonts)} 个拉丁字体面（中文 {F_CJK} 依赖系统字体，跨平台请见文档）")

# ──────────────────────────── 主流程 ────────────────────────────
def build(content_path, out_path):
    with open(content_path, encoding="utf-8") as fh:
        data = json.load(fh)
    meta = data.get("meta", {})
    global F_CJK
    F_CJK = meta.get("cjk_font", F_CJK)
    handle, site = meta.get("handle"), meta.get("site")
    base_dir = meta.get("base_image_dir")

    def resolve(p):
        if p and base_dir and not os.path.isabs(p): return os.path.join(base_dir, p)
        return p

    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H

    for i, s in enumerate(data.get("slides", [])):
        t = s.get("type")
        if t not in BUILDERS:
            print(f"  ⚠ 跳过未知类型: {t}"); continue
        lint(s, i)
        s.setdefault("_handle", handle); s.setdefault("_site", site)
        if s.get("image"): s["image"] = resolve(s["image"])
        if s.get("images"): s["images"] = [resolve(x) for x in s["images"]]
        sl = BUILDERS[t](prs, s)
        add_notes(sl, s.get("notes"))

    prs.save(out_path)
    if _warnings:
        print("内容 lint 警告（不阻断，供精简参考）:"); print("\n".join(_warnings))
    if meta.get("embed_fonts", True):
        embed_fonts(out_path)
    print(f"✅ 生成: {out_path}（{len(prs.slides._sldIdLst)} 页）")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("用法: python3 build_deck.py <content.json> <output.pptx>"); sys.exit(1)
    ensure_fonts()
    build(sys.argv[1], sys.argv[2])
